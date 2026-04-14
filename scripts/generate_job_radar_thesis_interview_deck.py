from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation

from generate_job_radar_interview_deck import (
    ACCENT_SOFT,
    BG,
    GOLD_SOFT,
    GREEN_SOFT,
    PANEL,
    RED_SOFT,
    SLIDE_H,
    SLIDE_W,
    add_arrow,
    add_badge,
    add_bullet_list,
    add_footer,
    add_flow_node,
    add_kpi,
    add_text_panel,
    add_title,
    set_background,
)


DEFAULT_OUTPUT = "/Users/zhuangcaizhen/Desktop/專案/UI練習/job-radar-thesis-interview-deck.pptx"


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    band = slide.shapes.add_shape(1, 0.7 * 914400, 0.62 * 914400, 12 * 914400, 5.95 * 914400)
    band.fill.solid()
    band.fill.fore_color.rgb = PANEL
    band.line.color.rgb = BG

    add_badge(slide, "THESIS × INTERVIEW", 0.95, 0.92, 2.15)
    add_title(
        slide,
        "Job Radar\n研究式講解 + 工作面試版簡報",
        "用論文報告的結構講清楚產品、系統、AI 與實驗，但保留一般工作面試最需要的商業價值與執行力。",
    )
    add_kpi(slide, 8.1, 1.7, 1.95, 1.22, "講法結構", "Problem", ACCENT_SOFT)
    add_kpi(slide, 10.2, 1.7, 1.95, 1.22, "講法結構", "Method", GOLD_SOFT)
    add_kpi(slide, 8.1, 3.05, 1.95, 1.22, "講法結構", "Result", GREEN_SOFT)
    add_kpi(slide, 10.2, 3.05, 1.95, 1.22, "講法結構", "Trade-off", RED_SOFT)
    add_text_panel(
        slide,
        0.98,
        3.08,
        6.0,
        2.15,
        "這份 deck 的目的",
        [
            "把專案從『我做了一個求職工具』提升成『我定義問題、設計系統、用 AI 做方法、用 benchmark 驗證、再持續優化』。",
            "面試時外層仍然是產品與工程敘事，內層才用研究式結構把方法與證據講紮實。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        7.25,
        4.48,
        5.0,
        1.45,
        "建議講法",
        [
            "10 分鐘版本：產品問題 1 分鐘、架構 2 分鐘、AI 方法 3 分鐘、實驗與優化 2 分鐘、擴充與 trade-off 2 分鐘。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "講解路徑", "這一頁是給面試官的導航，也讓你自己不會講散。")
    add_flow_node(slide, 0.9, 2.0, 2.15, 1.0, "Problem", "使用者痛點與產品定義", PANEL)
    add_flow_node(slide, 3.35, 2.0, 2.15, 1.0, "Method", "系統架構、資料流、AI 方法", ACCENT_SOFT)
    add_flow_node(slide, 5.8, 2.0, 2.15, 1.0, "Experiment", "benchmark、live latency、驗證方式", GOLD_SOFT)
    add_flow_node(slide, 8.25, 2.0, 2.15, 1.0, "Trade-off", "為什麼這樣設計、限制在哪", GREEN_SOFT)
    add_flow_node(slide, 10.7, 2.0, 1.65, 1.0, "Value", "你能帶來什麼", RED_SOFT)
    for start in [3.05, 5.5, 7.95, 10.4]:
        add_arrow(slide, start, 2.5, start + 0.15, 2.5)
    add_text_panel(
        slide,
        0.95,
        3.7,
        12.0,
        2.2,
        "實戰講法",
        [
            "如果面試官偏產品，就把時間放在 Problem、Value；如果偏工程或 AI，就把時間拉到 Method、Experiment。這份 deck 的設計就是讓你可以自由裁切，不會只能照順序背。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_contributions(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "我的角色與實際貢獻", "這一頁很重要，避免整場講完只剩系統很強、但你本人不清楚。")
    add_text_panel(
        slide,
        0.75,
        1.62,
        3.8,
        4.95,
        "我實際做的事",
        [
            "定義產品邊界，把專案從『資料抓取工具』提升成『求職工作台』。",
            "拆系統層次：UI、runtime、pipeline、AI、DB、history。",
            "設計並優化 RAG：chunking、retrieval、rerank、ANN、latency。",
            "做大規模重構，把大檔拆成 page / sections / actions / helpers。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.85,
        1.62,
        3.75,
        4.95,
        "我怎麼證明不是空談",
        [
            "有 benchmark、live metrics、testing 手冊、AI / RAG 報告頁。",
            "有具體 refactor 輪次與模組拆分結果。",
            "有 production 設定、maintenance、cache 治理、canonical normalization。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.9,
        1.62,
        3.6,
        4.95,
        "面試時的說法",
        [
            "我不是只做某個功能，而是從問題定義、架構拆分、AI 方法、驗證到持續優化都真的碰過。",
            "這代表我能 handle 一個完整產品，而不是只寫單點程式。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_role_fit(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "這個職缺我會怎麼對位", "不把重點放在特定雲商，而是把可遷移能力講清楚。")
    add_text_panel(
        slide,
        0.75,
        1.62,
        3.85,
        4.95,
        "已經具備",
        [
            "LLM / RAG / AI 助理設計與落地。",
            "多來源資料整合、正規化、分析與通知流程。",
            "把 AI 功能做成可使用、可維運的產品。",
            "系統重構、模組邊界整理、文件化輸出。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.83,
        1.62,
        3.85,
        4.95,
        "可直接對位到 JD",
        [
            "AI 解決方案導入與客製化。",
            "對話式 AI 應用建構。",
            "技術支援、問題排除與穩定運行。",
            "需求分析、團隊溝通與文件報告。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.92,
        1.62,
        3.6,
        4.95,
        "我會怎麼講",
        [
            "我不會硬說自己精通某個特定雲平台，而是誠實講：我已經有把 AI 能力整合成產品、資料流與維運流程的經驗，這些能力本身就是最重要的底層能力。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_thesis_framing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "1. 為什麼這個專案適合用論文式結構來講", "因為它不只是產品功能，而是有清楚的問題、方法、資料流與驗證。")
    add_text_panel(
        slide,
        0.75,
        1.62,
        3.8,
        4.95,
        "研究式結構",
        [
            "Problem：求職資訊分散，使用者缺乏持續回訪與決策支持。",
            "Method：建立 snapshot-centered 工作台，結合 crawler、analysis、RAG、履歷匹配與通知機制。",
            "Evaluation：用 retrieval benchmark、live latency、系統資料分層來驗證系統品質。",
            "Future Work：向量索引、chunking、canonical normalization、產品擴充。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.85,
        1.62,
        3.75,
        4.95,
        "面試官真正會聽什麼",
        [
            "你能不能把一個模糊題目拆成明確問題。",
            "你有沒有方法，不是只靠直覺做功能。",
            "你有沒有用數據或實驗驗證，而不是『我覺得比較好』。",
            "你知不知道系統限制，能不能講出下一步。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.9,
        1.62,
        3.6,
        4.95,
        "你要避免的事",
        [
            "不要把它講成學術口試，面試官不是來聽文獻回顧。",
            "不要一開始就陷進 chunking / rerank 細節，先講產品價值。",
            "不要只講技術漂亮，忽略使用者場景與為什麼這樣拆。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_user_loop(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "使用者工作流與回訪循環", "這一頁補上產品視角，避免簡報變成只有系統圖。")
    add_flow_node(slide, 0.85, 2.02, 2.05, 0.95, "輸入條件", "角色 / 地區 / 關鍵字 / 頻率", PANEL)
    add_flow_node(slide, 3.18, 2.02, 2.05, 0.95, "抓取快照", "搜尋結果、分析結果、最新變化", ACCENT_SOFT)
    add_flow_node(slide, 5.51, 2.02, 2.05, 0.95, "做決策", "收藏、履歷匹配、AI 問答", GOLD_SOFT)
    add_flow_node(slide, 7.84, 2.02, 2.05, 0.95, "持續追蹤", "saved search、看板、通知", GREEN_SOFT)
    add_flow_node(slide, 10.17, 2.02, 2.05, 0.95, "再次回訪", "下一輪刷新與投遞", RED_SOFT)
    for start in [2.95, 5.28, 7.61, 9.94]:
        add_arrow(slide, start, 2.47, start + 0.15, 2.47)
    add_text_panel(
        slide,
        0.9,
        3.65,
        12.0,
        2.3,
        "這頁的價值",
        [
            "面試時你可以直接說：這個產品不是查完就結束，而是要形成回訪循環。這就是為什麼系統後面一定會長出 saved search、自動刷新、tracking center、notifications，因為這些不是附加功能，而是產品閉環的一部分。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_problem_statement(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "2. 問題定義：這個系統到底要解什麼", "先把研究問題講清楚，後面的方法與架構才有依據。")
    add_text_panel(
        slide,
        0.75,
        1.55,
        4.0,
        5.0,
        "核心問題",
        [
            "職缺分散在多個平台，使用者要重複搜尋、重複判斷、重複整理。",
            "市場資訊是動態的，但大部分求職工具沒有長期回訪與自動刷新能力。",
            "履歷與市場需求之間缺少可量化的對照機制。",
            "如果把 AI 接進來，回答又不能脫離實際市場資料。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        5.0,
        1.55,
        3.7,
        5.0,
        "因此我定義的產品目標",
        [
            "把一次性的查詢流程，變成可持續回訪的求職工作台。",
            "把原始職缺資料轉成 snapshot、insight、匹配結果與通知信號。",
            "讓 AI 的回答建立在可追溯的市場上下文，而不是純聊天。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.55,
        3.55,
        5.0,
        "面試時一句話講法",
        [
            "我不是在做職缺爬蟲，而是在做一個以市場快照為核心的求職工作台，讓搜尋、分析、決策、回訪與 AI 問答形成閉環。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_research_questions(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "3. 研究問題與系統問題", "這一頁就是論文式思考最有感的地方。")
    add_bullet_list(
        slide,
        0.95,
        1.75,
        11.0,
        2.3,
        [
            "RQ1：如何把多來源職缺資料轉成穩定、可比較、可持續刷新的市場快照？",
            "RQ2：如何在不依賴純黑箱模型的前提下，做出可解釋的履歷匹配與市場分析？",
            "RQ3：如何讓 AI 助理大多數問題都能 grounded 到市場資料，同時保留對一般問題的自然互動能力？",
            "RQ4：在真實產品約束下，如何平衡 retrieval 精準度、平均延遲與系統複雜度？",
        ],
        font_size=17,
    )
    add_text_panel(
        slide,
        0.95,
        4.45,
        12.0,
        1.75,
        "面試時的價值",
        [
            "這樣講會讓面試官知道，你不是一邊做一邊堆功能，而是有明確 problem framing。工程設計、AI 方法、資料分層，都是從這些問題往下推導出來的。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "4. 方法總覽：整體系統架構", "先把方法講成一張可理解的系統圖。")
    add_flow_node(slide, 0.85, 2.0, 2.0, 1.05, "Frontend", "Streamlit UI / session / router", PANEL)
    add_flow_node(slide, 3.25, 2.0, 2.1, 1.05, "Application", "bootstrap / crawl service / runtime", ACCENT_SOFT)
    add_flow_node(slide, 5.8, 2.0, 2.0, 1.05, "Pipeline", "connector / normalize / dedupe / analyze", GOLD_SOFT)
    add_flow_node(slide, 8.2, 2.0, 2.0, 1.05, "AI Layer", "resume analysis / RAG / monitoring", GREEN_SOFT)
    add_flow_node(slide, 10.6, 2.0, 1.75, 1.05, "Storage", "SQLite / cache / history", RED_SOFT)
    for start in [2.95, 5.5, 7.95, 10.35]:
        add_arrow(slide, start, 2.55, start + 0.18, 2.55)
    add_text_panel(
        slide,
        0.95,
        3.7,
        12.0,
        2.25,
        "我會怎麼講這張圖",
        [
            "這是一個模組化單體。前端負責工作台互動；application layer 協調同步與非同步任務；pipeline 把原始資料轉成可用快照；AI layer 建立履歷匹配與市場問答；storage 則用多庫分離產品主資料、執行期資料與歷史資料。",
            "這樣切的好處是邏輯邊界清楚，壞處是邊界多了之後需要更嚴格的資料契約與觀測。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_data_method(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "5. 資料方法：從原始職缺到市場快照", "這一頁把資料工程與特徵工程講清楚。")
    add_flow_node(slide, 0.8, 2.0, 2.0, 0.95, "Raw Jobs", "104 / 1111 / Cake / LinkedIn", PANEL)
    add_flow_node(slide, 3.15, 2.0, 2.0, 0.95, "Cleaning", "canonical company / title / location", ACCENT_SOFT)
    add_flow_node(slide, 5.5, 2.0, 2.0, 0.95, "Merge", "cross-source dedupe / enrich", GOLD_SOFT)
    add_flow_node(slide, 7.85, 2.0, 2.0, 0.95, "Analysis", "role / skills / tasks / salary", GREEN_SOFT)
    add_flow_node(slide, 10.2, 2.0, 2.0, 0.95, "Snapshot", "market view / AI context", RED_SOFT)
    for start in [2.9, 5.25, 7.6, 9.95]:
        add_arrow(slide, start, 2.48, start + 0.15, 2.48)
    add_text_panel(
        slide,
        0.9,
        3.65,
        5.85,
        2.35,
        "為什麼這樣做",
        [
            "如果直接把 raw jobs 丟給前端或 AI，資料會非常噪、重複且難比較。",
            "所以我先做 canonical normalization、跨來源同職缺合併、技能/任務/角色萃取，讓後面的 UI、匹配與 RAG 都建立在結構化資料上。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.0,
        3.65,
        5.45,
        2.35,
        "面試時可延伸",
        [
            "這其實就是 feature engineering，只是它不是單純為了模型，而是同時服務產品顯示、排序、匹配、檢索與 dedupe。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_db_boundary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "資料與 DB 邊界", "多資料庫不是炫技，而是 ownership 與生命週期不同。")
    add_text_panel(
        slide,
        0.75,
        1.58,
        2.9,
        5.05,
        "product_state",
        [
            "產品主資料。",
            "使用者偏好、saved search、收藏、通知設定、AI monitoring 等。",
            "是最接近產品面和 UI 面的資料庫。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        3.88,
        1.58,
        2.9,
        5.05,
        "query_runtime",
        [
            "執行期資料。",
            "crawl_jobs、signals、runtime queue、job status。",
            "它的生命週期短，主要服務 scheduler / worker / UI 輪詢。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        7.01,
        1.58,
        2.9,
        5.05,
        "market_history",
        [
            "歷史資料。",
            "保留長期職缺 archive，未來可支撐趨勢分析與全庫檢索。",
            "跟當前 snapshot 的用途不同。",
        ],
        fill=GOLD_SOFT,
    )
    add_text_panel(
        slide,
        10.14,
        1.58,
        2.35,
        5.05,
        "cache / snapshots",
        [
            "非正式主資料，但對效能很重要。",
            "服務 embeddings、RAG、最新快照與暫存結果。",
            "需要 maintenance，不然會長髒。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_delivery_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "導入與交付流程", "這一頁補上『我如何把技術變成可交付方案』。")
    add_flow_node(slide, 1.0, 2.1, 2.2, 0.95, "需求訪談", "場景、資料、限制、成功條件", PANEL)
    add_flow_node(slide, 3.6, 2.1, 2.2, 0.95, "PoC 定義", "要先證明什麼、如何驗證", ACCENT_SOFT)
    add_flow_node(slide, 6.2, 2.1, 2.2, 0.95, "導入實作", "資料流、AI 流程、介面、權限", GREEN_SOFT)
    add_flow_node(slide, 8.8, 2.1, 2.2, 0.95, "驗收與交付", "文件、FAQ、操作說明、風險", GOLD_SOFT)
    add_arrow(slide, 3.3, 2.55, 3.48, 2.55)
    add_arrow(slide, 5.9, 2.55, 6.08, 2.55)
    add_arrow(slide, 8.5, 2.55, 8.68, 2.55)
    add_text_panel(
        slide,
        1.0,
        3.75,
        11.2,
        2.25,
        "面試時這樣講",
        [
            "我的強項不是只把功能做出來，而是能把需求拆成導入流程：先定義問題與成功指標，再做 PoC，然後把系統、AI、資料與文件一起交付。這個專案雖然是個人產品，但我已經把這條思路走過一次。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_ai_method(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "6. AI 方法：LLM / RAG / Matching 怎麼放進產品", "重點不是有沒有接模型，而是模型在整體系統裡扮演什麼角色。")
    add_text_panel(
        slide,
        0.75,
        1.58,
        3.9,
        5.0,
        "履歷分析",
        [
            "先做 rule-based 萃取與 context selection，再用 OpenAI 補摘要與 profile 結構化。",
            "匹配時混合 exact overlap、keyword、semantic similarity 與 role/title signal。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.9,
        1.58,
        3.75,
        5.0,
        "RAG 問答",
        [
            "以 snapshot 為中心建立 chunks，透過 hybrid retrieval 找出相關市場證據，再由 LLM 生成 grounded 回答。",
            "大多數帶時間脈絡或求職脈絡的問題都會走 retrieval，不只侷限於預設職涯題。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.9,
        1.58,
        3.6,
        5.0,
        "可靠性設計",
        [
            "AI 不是 source of truth，而是建立在 market snapshot 的檢索上下文。",
            "有 monitoring、mode routing、fallback、latency 優化與 benchmark。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_llm_contract(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "LLM 合約：模型不是隨便接上去", "這一頁補上 prompt / mode / guardrail 思維，會更像成熟 AI 系統。")
    add_text_panel(
        slide,
        0.75,
        1.62,
        3.9,
        4.95,
        "mode routing",
        [
            "問題先分 answer mode，再決定 retrieval policy、top-k、max output tokens。",
            "求職題走較正式 QA；一般題保留自然互動，但大多數帶時間點的題仍會進 retrieval。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.9,
        1.62,
        3.75,
        4.95,
        "grounding contract",
        [
            "不是所有問題都讓模型自由發揮，而是要求它優先依據 snapshot context 回答。",
            "這樣能讓答案更像『市場顧問』，而不是泛知識聊天機器人。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.9,
        1.62,
        3.6,
        4.95,
        "輸出控制",
        [
            "不同 mode 有不同 token budget，目的是控制平均延遲與回答密度。",
            "這其實是產品決策，不只是模型參數調整。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_rag_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "7. RAG 流程：這是我最像研究方法的一段", "把 chunking、retrieval、rerank、ANN 與延遲優化一次講清楚。")
    add_flow_node(slide, 0.7, 2.02, 1.65, 0.9, "Question", "使用者提問", PANEL)
    add_flow_node(slide, 2.65, 2.02, 1.8, 0.9, "Mode / Intent", "分類題型與 retrieval policy", ACCENT_SOFT)
    add_flow_node(slide, 4.75, 2.02, 1.8, 0.9, "Chunk Build", "snapshot / resume / market chunks", GOLD_SOFT)
    add_flow_node(slide, 6.85, 2.02, 1.8, 0.9, "ANN Recall", "persistent index 候選召回", GREEN_SOFT)
    add_flow_node(slide, 8.95, 2.02, 1.8, 0.9, "Hybrid Rerank", "embedding + lexical + signal", ACCENT_SOFT)
    add_flow_node(slide, 11.05, 2.02, 1.55, 0.9, "Answer", "LLM grounded response", RED_SOFT)
    for start in [2.42, 4.52, 6.62, 8.72, 10.82]:
        add_arrow(slide, start, 2.47, start + 0.14, 2.47)
    add_text_panel(
        slide,
        0.85,
        3.45,
        12.0,
        2.45,
        "我做的優化",
        [
            "chunking 從結構化摘要進一步提升到 anchored chunk 與 item-level chunks；retrieval 上加入 job-specific bonus、market penalty、動態 top-k、comparison coverage；再往上加長期持久化 ANN 候選召回。",
            "延遲上則用 chunk cache、embedding memory cache、snapshot hash sync、mode-based output budget 來壓平均反應時間。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_experiment(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "8. 實驗設計：我怎麼驗證系統不是只靠感覺", "把 benchmark、live metrics、manual regression 放在一起看。")
    add_text_panel(
        slide,
        0.78,
        1.6,
        3.95,
        4.95,
        "離線 benchmark",
        [
            "用 snapshot 自動生成 default 與 realistic evaluation set。",
            "比較不同 chunking strategy、top-k、rerank 規則與 retrieval coverage。",
            "指標看 hit@1、hit@3、hit@5、MRR。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.93,
        1.6,
        3.85,
        4.95,
        "live metrics",
        [
            "AI monitoring event 會記錄 answer mode、latency、status 與 usage。",
            "用 live DB 看 market_summary、general_chat、resume.build_profile 等真實平均延遲。",
            "這能把『程式設計』和『真實上線狀態』分開看。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.98,
        1.6,
        3.55,
        4.95,
        "手動回歸",
        [
            "重要 UI 流程仍需要人工驗證：查詢、saved search、通知、AI 助理、履歷分析。",
            "因為最終產品不是只看 retrieval 分數，而是整條使用者工作流是否穩定。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_benchmark_detail(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Benchmark 細節：我怎麼設計 evaluation set", "這一頁會讓你看起來不像只會調參，而是真的在做 AI 工程驗證。")
    add_text_panel(
        slide,
        0.78,
        1.62,
        3.95,
        4.95,
        "default set",
        [
            "偏系統型題目。",
            "包含市場技能、工作內容、來源分布、角色分布、薪資等問題。",
            "主要用來檢查 retrieval 主幹有沒有壞。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.93,
        1.62,
        3.85,
        4.95,
        "realistic set",
        [
            "更像真實使用者提問。",
            "例如『如果先看整體市場...』『如果投這份工作...』。",
            "目前 realistic set 更接近線上真實壓力。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.98,
        1.62,
        3.55,
        4.95,
        "我怎麼用它",
        [
            "先比 chunking，再比 retrieval rule，再比 top-k 與 rerank。",
            "不是直接改 production，而是先在 benchmark 確認方向，再往正式流程接。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_results(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "9. 結果與目前證據", "這一頁是你面試時最有說服力的數據頁。")
    add_kpi(slide, 0.92, 1.85, 2.25, 1.25, "Retrieval realistic", "Hit@1 78.95%", ACCENT_SOFT)
    add_kpi(slide, 3.4, 1.85, 2.25, 1.25, "Retrieval realistic", "MRR 0.807", GOLD_SOFT)
    add_kpi(slide, 5.88, 1.85, 2.25, 1.25, "Assistant latency", "7021.93 ms", GREEN_SOFT)
    add_kpi(slide, 8.36, 1.85, 2.25, 1.25, "Resume build_profile", "5973.97 ms", RED_SOFT)
    add_kpi(slide, 10.84, 1.85, 1.55, 1.25, "Cache", "255.9 MB", ACCENT_SOFT)
    add_text_panel(
        slide,
        0.9,
        3.45,
        5.85,
        2.4,
        "可以怎麼講",
        [
            "我不是只說『AI 助理有變好』，而是有 benchmark 能看 retrieval accuracy，也有 live metrics 能看真實 latency。",
            "這代表我把 AI 當成工程系統來做，而不是只停在 demo 或 prompt 試玩。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.0,
        3.45,
        5.45,
        2.4,
        "也要誠實講限制",
        [
            "目前 retrieval benchmark 仍是離線代理 embedding，不等於線上最終回答品質。",
            "持久化 ANN index 已接上，但 live corpus 還需要更多歷史資料才會發揮全值。",
            "平均延遲有下降空間，最大頭仍是模型網路往返。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_optimization_path(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "優化歷程：不是一次到位，而是逐步收斂", "這一頁能讓你講出工程判斷與迭代能力。")
    add_flow_node(slide, 0.8, 2.0, 2.2, 1.0, "Stage 1", "structured chunks + 基本 hybrid retrieval", PANEL)
    add_flow_node(slide, 3.35, 2.0, 2.2, 1.0, "Stage 2", "anchored chunks + job-specific rerank", ACCENT_SOFT)
    add_flow_node(slide, 5.9, 2.0, 2.2, 1.0, "Stage 3", "dynamic top-k + market penalty", GOLD_SOFT)
    add_flow_node(slide, 8.45, 2.0, 2.2, 1.0, "Stage 4", "persistent ANN + cache optimizations", GREEN_SOFT)
    add_flow_node(slide, 11.0, 2.0, 1.55, 1.0, "Next", "answer-level eval / richer corpus", RED_SOFT)
    for start in [3.1, 5.65, 8.2, 10.75]:
        add_arrow(slide, start, 2.5, start + 0.15, 2.5)
    add_text_panel(
        slide,
        0.95,
        3.7,
        12.0,
        2.2,
        "這一頁的面試意義",
        [
            "你可以很清楚地講出：我不是直接丟一個 AI 架構上去，而是先建立 baseline，再逐步修 chunking、rerank、top-k、ANN、latency。這會讓面試官看到你是會做 experiment loop 的工程師。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_ops_cases(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "維運與問題排除案例", "這一頁直接對應職缺裡『請準備維運問題案例說明』。")
    add_text_panel(
        slide,
        0.75,
        1.62,
        3.85,
        4.95,
        "案例一：AI 助理延遲",
        [
            "問題：回應平均延遲偏高，使用體驗不穩。",
            "處理：chunk cache、embedding memory cache、snapshot hash sync、mode-based token budget。",
            "價值：不是只會調 prompt，而是會拆熱路徑與重複工作。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.85,
        1.62,
        3.85,
        4.95,
        "案例二：RAG 精準度",
        [
            "問題：特定職缺問題容易被市場級 chunk 搶走。",
            "處理：anchored chunk、job-specific rerank、top-k routing、market penalty、evaluation set。",
            "價值：把『AI 不準』拆成可量測、可優化的檢索問題。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.62,
        3.55,
        4.95,
        "案例三：cache 維運",
        [
            "問題：cache 超過設定上限，開始污染磁碟與效能。",
            "處理：固定 maintenance、live cleanup、上限控制、統計回報。",
            "價值：證明我有持續運營與維護系統的意識。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_tradeoffs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "10. Trade-off：我為什麼這樣設計，而不是別種做法")
    add_text_panel(
        slide,
        0.75,
        1.65,
        3.9,
        4.9,
        "模組化單體 vs 微服務",
        [
            "我選模組化單體，因為 stage 還沒到要為網路邊界與服務編排付出更高複雜度。",
            "優點是快、可控；缺點是需要更嚴格的模組邊界與重構紀律。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.9,
        1.65,
        3.8,
        4.9,
        "snapshot-centered RAG vs 全域知識庫",
        [
            "先以當次市場快照為核心，能保證答案比較貼近『目前這批資料』。",
            "缺點是歷史與跨 corpus 問題能力較弱，所以後續才補 persistent ANN 與全庫召回。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.65,
        3.55,
        4.9,
        "規則 + AI 混合",
        [
            "我沒有把所有決策都交給模型，而是先用 canonicalization、features、rerank rule 穩住系統。",
            "優點是可控且可解釋；缺點是規則會累積，需要定期收斂。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_documentation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "文件與報告能力", "這個 JD 很重技術文件與知識分享，這是你的加分題。")
    add_text_panel(
        slide,
        0.78,
        1.62,
        3.95,
        4.95,
        "我做過的文件",
        [
            "系統架構手冊。",
            "AI / LLM / RAG 詳細報告。",
            "DB schema 手冊。",
            "testing 手冊。",
            "產品與系統設計 playbook。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.93,
        1.62,
        3.85,
        4.95,
        "它們服務誰",
        [
            "工程師：看架構、DB、RAG、測試。",
            "管理者：看系統能力、風險、路線圖。",
            "面試官 / 利害關係人：快速理解你做的是什麼。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.98,
        1.62,
        3.55,
        4.95,
        "面試時怎麼講",
        [
            "我不只會做系統，還會把系統講清楚、寫清楚、交接清楚。這對導入型職缺特別重要，因為技術能不能被持續使用，很大一部分取決於文件與溝通品質。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_why_me(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "為什麼這個專案能代表我的價值", "最後要把系統收回到你本人。")
    add_text_panel(
        slide,
        0.8,
        1.65,
        5.9,
        4.95,
        "這個專案證明的能力",
        [
            "我能從需求與場景出發定義產品，而不是只追功能。",
            "我能把複雜系統拆層、重構、維持可演進性。",
            "我能把 AI 放在正確位置：有資料邊界、有檢索、有監控、有 trade-off。",
            "我能用 benchmark 與 live metrics 驗證，而不是只說『我覺得比較好』。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        6.95,
        1.65,
        5.55,
        4.95,
        "一句話收尾",
        [
            "如果把這個專案當成作品，它證明我不只會寫程式，而是能把問題定義、系統設計、AI 方法、驗證與產品思考整合成一個完整可講清楚的成果。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_future_work(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "11. Future Work：如果要繼續研究或產品化，我會往哪裡走")
    add_bullet_list(
        slide,
        0.95,
        1.72,
        11.1,
        2.7,
        [
            "把 evaluation set 擴成更多真實使用者多輪提問，並把 answer-level quality 也納進評估。",
            "把 chunking 從結構化摘要進一步推到 paragraph-aware / recursive split 的 job description segmentation。",
            "讓 persistent vector index 真正吃到更多歷史 snapshot 與多 corpus metadata filter。",
            "把 AI 成本、延遲與 quality 做成長期儀表板，形成 steady-state 對照。",
        ],
        font_size=17,
    )
    add_text_panel(
        slide,
        0.95,
        4.75,
        12.0,
        1.45,
        "面試時一句話",
        [
            "我不是只把專案做到能用，而是已經知道下一輪研究與工程演進要往哪裡走，這代表我對系統邊界和成長路線有清楚判斷。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_talk_track(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "12. 面試實戰講法：同一份內容，怎麼講得像研究生也像工程師")
    add_text_panel(
        slide,
        0.82,
        1.65,
        6.0,
        4.95,
        "5 分鐘版",
        [
            "先用 1 分鐘講產品問題與定位。",
            "再用 1 分鐘講整體架構。",
            "接著用 1.5 分鐘講 AI / RAG 方法。",
            "再用 1 分鐘講 benchmark 與 live latency。",
            "最後 30 秒講 trade-off 與未來方向。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.05,
        1.65,
        5.45,
        4.95,
        "回答追問的原則",
        [
            "追問技術時：講方法、數據、限制。",
            "追問產品時：講使用者工作流、回訪、留存與決策支持。",
            "追問 AI 時：講 grounded context、benchmark、latency、trade-off。",
            "追問你本人時：講你定義問題、做拆分、做重構、做驗證。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def slide_appendix(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Appendix. 建議一起帶去面試的材料")
    add_bullet_list(
        slide,
        0.95,
        1.75,
        11.2,
        4.9,
        [
            "這份 thesis-style 面試 deck。",
            "通用版系統 deck：job-radar-interview-deck.pptx。",
            "這次應徵職缺的客製化說法：導入、維運、文件、AI 助理案例。",
            "AI / LLM / RAG 報告頁與完整結構圖頁。",
            "如果有時間，再準備一張白板版：畫 app / pipeline / RAG / DB 的總圖。",
        ],
        font_size=17,
    )
    add_footer(slide, "Job Radar · Thesis-Style Interview Deck")


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_agenda(prs)
    slide_contributions(prs)
    slide_role_fit(prs)
    slide_thesis_framing(prs)
    slide_user_loop(prs)
    slide_problem_statement(prs)
    slide_research_questions(prs)
    slide_architecture(prs)
    slide_data_method(prs)
    slide_db_boundary(prs)
    slide_delivery_flow(prs)
    slide_ai_method(prs)
    slide_llm_contract(prs)
    slide_rag_pipeline(prs)
    slide_experiment(prs)
    slide_benchmark_detail(prs)
    slide_results(prs)
    slide_optimization_path(prs)
    slide_ops_cases(prs)
    slide_tradeoffs(prs)
    slide_documentation(prs)
    slide_why_me(prs)
    slide_future_work(prs)
    slide_talk_track(prs)
    slide_appendix(prs)
    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate Job Radar thesis-style interview deck.")
    parser.add_argument("--output", default=DEFAULT_OUTPUT, help="Output PPTX path")
    args = parser.parse_args()

    output = Path(args.output).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)
    build_presentation().save(output)
    print(output)


if __name__ == "__main__":
    main()
