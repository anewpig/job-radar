from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation

from generate_job_radar_interview_deck import (
    ACCENT_SOFT,
    BG,
    GOLD_SOFT,
    GREEN_SOFT,
    PANEL,
    RED_SOFT,
    SLIDE_H,
    SLIDE_W,
    add_arrow,
    add_badge,
    add_bullet_list,
    add_footer,
    add_flow_node,
    add_kpi,
    add_text_panel,
    add_title,
    set_background,
)


DEFAULT_OUTPUT = "/Users/zhuangcaizhen/Desktop/專案/UI練習/job-radar-helixgpt-agent-deck.pptx"


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    band = slide.shapes.add_shape(1, 0.7 * 914400, 0.62 * 914400, 12 * 914400, 5.95 * 914400)
    band.fill.solid()
    band.fill.fore_color.rgb = PANEL
    band.line.color.rgb = BG
    add_badge(slide, "HELIXGPT × AGENT FIT", 0.95, 0.92, 2.35)
    add_title(
        slide,
        "Job Radar\nhelixGPT 類 AI Agent 職缺對位版",
        "把你的專案改寫成『企業 AI 助手 / AI agent / 知識導入』語言，方便你直接對面試官講。",
    )
    add_kpi(slide, 8.0, 1.7, 1.95, 1.22, "角色重點", "Agent", ACCENT_SOFT)
    add_kpi(slide, 10.1, 1.7, 1.95, 1.22, "角色重點", "RAG", GOLD_SOFT)
    add_kpi(slide, 8.0, 3.05, 1.95, 1.22, "角色重點", "Ops", GREEN_SOFT)
    add_kpi(slide, 10.1, 3.05, 1.95, 1.22, "角色重點", "Enablement", RED_SOFT)
    add_text_panel(
        slide,
        0.98,
        3.08,
        6.05,
        2.15,
        "這份 deck 的主張",
        [
            "如果這個職缺實際是在做 helixGPT 類企業 AI agent，那我最該證明的不是某個特定雲商，而是：我能把資料、檢索、工具、回答、維運與文件整成可交付的 AI 助手產品。",
            "這份簡報就是沿著這個方向重寫。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_what_role_is(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "1. 如果這個職缺在做 helixGPT，本質上是在做什麼")
    add_text_panel(
        slide,
        0.75,
        1.6,
        3.9,
        4.95,
        "我對角色的理解",
        [
            "企業 AI 助手，不只是聊天介面。",
            "要處理知識檢索、工作流、權限、穩定性、工具整合與使用者導入。",
            "核心價值不是模型本身，而是把 AI 能力嵌進使用者工作流程。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.9,
        1.6,
        3.75,
        4.95,
        "這個職缺真正要看",
        [
            "你會不會做對話式 AI 與企業知識問答。",
            "你能不能把需求拆成 AI workflow。",
            "你能不能做導入、穩定運行與技術說明。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.9,
        1.6,
        3.6,
        4.95,
        "面試時一句話",
        [
            "如果把 helixGPT 看成企業知識助手或 AI agent 平台，那我已經做過最接近的事情，就是把資料、檢索、分析、問答與長期維運整成一個可用產品。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_mapping(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "2. Job Radar 如何映射到 helixGPT 類產品")
    add_text_panel(
        slide,
        0.78,
        1.58,
        3.95,
        5.0,
        "Job Radar 現有能力",
        [
            "市場快照與資料清理。",
            "RAG 問答與 grounded 回答。",
            "履歷分析與個人化建議。",
            "saved search、自動刷新、通知與 tracking。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.93,
        1.58,
        3.85,
        5.0,
        "對應 helixGPT 類場景",
        [
            "企業知識庫問答。",
            "AI 助手 / 企業智能客服。",
            "工作台式回訪與持續提醒。",
            "以角色與情境為中心的個人化回答。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.98,
        1.58,
        3.55,
        5.0,
        "可轉移能力",
        [
            "我不是只做求職場景，而是已經做過一套『資料 -> 檢索 -> 回答 -> 回訪』的完整 AI 助手閉環。",
            "場景換成企業知識，底層方法仍然成立。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_agent_arch(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "3. 我會怎麼講 AI Agent 架構")
    add_flow_node(slide, 0.85, 2.05, 1.8, 0.95, "User Query", "使用者問題 / 任務", PANEL)
    add_flow_node(slide, 2.95, 2.05, 1.95, 0.95, "Intent Router", "問題類型 / 模式 / 策略", ACCENT_SOFT)
    add_flow_node(slide, 5.2, 2.05, 1.9, 0.95, "Knowledge Layer", "snapshot / chunks / retrieval", GOLD_SOFT)
    add_flow_node(slide, 7.4, 2.05, 1.9, 0.95, "Tool / Action", "查詢、分析、推薦、通知", GREEN_SOFT)
    add_flow_node(slide, 9.6, 2.05, 1.9, 0.95, "Response Layer", "grounded answer / next action", ACCENT_SOFT)
    add_flow_node(slide, 11.8, 2.05, 0.95, 0.95, "Observe", "log / metrics", RED_SOFT)
    for start in [2.7, 4.95, 7.15, 9.35, 11.55]:
        add_arrow(slide, start, 2.5, start + 0.15, 2.5)
    add_text_panel(
        slide,
        0.95,
        3.7,
        12.0,
        2.25,
        "我不會把 agent 講得太玄",
        [
            "我會直接說：我現在這套系統還不是複雜 planner 型 agent，但已經有 intent routing、knowledge grounding、tool-like workflow、monitoring 與回訪機制。對企業 AI 助手產品來說，這其實比炫技式多 agent 更重要。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_use_cases(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "4. 企業 AI 助手常見場景，我可以怎麼對位")
    add_text_panel(
        slide,
        0.75,
        1.6,
        3.85,
        4.95,
        "知識問答",
        [
            "把資料整理成可檢索 chunks。",
            "讓模型回答時有 grounded context。",
            "避免純泛知識亂答。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.85,
        1.6,
        3.85,
        4.95,
        "智能助手",
        [
            "根據使用者角色、問題脈絡與歷史狀態做個人化回覆。",
            "不只是回答，還能引導下一步動作。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.6,
        3.55,
        4.95,
        "流程型支援",
        [
            "用通知、saved search、tracking 這類概念，對位成企業裡的 follow-up、待辦提醒、狀態回訪。",
            "這也是 agent 產品常見的價值來源。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_delivery(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "5. 如果這是 AI agent 導入職缺，我會怎麼帶專案")
    add_flow_node(slide, 1.0, 2.1, 2.2, 0.95, "需求盤點", "使用者、知識來源、目標場景", PANEL)
    add_flow_node(slide, 3.6, 2.1, 2.2, 0.95, "PoC 定義", "先證明回答品質與使用價值", ACCENT_SOFT)
    add_flow_node(slide, 6.2, 2.1, 2.2, 0.95, "產品實作", "檢索、回答、工具、權限", GREEN_SOFT)
    add_flow_node(slide, 8.8, 2.1, 2.2, 0.95, "上線維運", "監控、問題排除、優化", GOLD_SOFT)
    add_arrow(slide, 3.3, 2.55, 3.48, 2.55)
    add_arrow(slide, 5.9, 2.55, 6.08, 2.55)
    add_arrow(slide, 8.5, 2.55, 8.68, 2.55)
    add_text_panel(
        slide,
        1.0,
        3.75,
        11.2,
        2.25,
        "這頁的說法",
        [
            "我會強調：我擅長的不是某個名詞，而是把 AI 助手從概念一路走到交付。這個專案雖然是求職場景，但我已經把資料、AI、UI、維運和文件一起做過，所以切到企業助理場景是合理的。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_ops(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "6. 維運案例：這會讓面試官覺得你真的碰過系統")
    add_text_panel(
        slide,
        0.75,
        1.62,
        3.85,
        4.95,
        "精準度問題",
        [
            "特定職缺問題容易抓到市場級資訊。",
            "我用 chunking、top-k routing、rerank、evaluation set 去調整。",
            "這很像企業 agent 裡常見的『答非所問』問題。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.85,
        1.62,
        3.85,
        4.95,
        "延遲問題",
        [
            "我不是只調 prompt，而是處理 chunk cache、embedding cache、snapshot sync、output budget。",
            "這表示我有 AI 系統的 performance engineering 意識。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.62,
        3.55,
        4.95,
        "維護問題",
        [
            "cache maintenance、canonical normalization、cross-source merge、監控事件。",
            "這些都會直接對位到企業產品上線後的日常維運。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_metrics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "7. 我不是只會講故事，還有數據")
    add_kpi(slide, 0.95, 1.9, 2.25, 1.25, "RAG realistic", "Hit@1 78.95%", ACCENT_SOFT)
    add_kpi(slide, 3.45, 1.9, 2.25, 1.25, "RAG realistic", "MRR 0.807", GOLD_SOFT)
    add_kpi(slide, 5.95, 1.9, 2.25, 1.25, "Assistant", "7021.93 ms", GREEN_SOFT)
    add_kpi(slide, 8.45, 1.9, 2.25, 1.25, "build_profile", "5973.97 ms", RED_SOFT)
    add_kpi(slide, 10.95, 1.9, 1.45, 1.25, "Events", "42+", ACCENT_SOFT)
    add_text_panel(
        slide,
        0.95,
        3.55,
        12.0,
        2.3,
        "面試時怎麼講",
        [
            "我會說：這套 AI 助手不是只是接模型而已，我有 benchmark 看 retrieval，也有 live metrics 看延遲。這樣面試官會直接把你歸類到『會做 AI 系統』，而不是『做過幾個 prompt demo』。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_docs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "8. 文件與 enablement：這很像內部 AI 產品推廣")
    add_text_panel(
        slide,
        0.78,
        1.62,
        3.95,
        4.95,
        "我已經做的文件",
        [
            "系統架構手冊。",
            "AI / LLM / RAG 報告。",
            "DB schema 手冊。",
            "testing 手冊。",
            "產品與學習手冊。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.93,
        1.62,
        3.85,
        4.95,
        "如果是企業 AI 產品",
        [
            "這些能力可以轉成 admin guide、user guide、FAQ、導入簡報、維運手冊。",
            "對內部 AI 助手來說，enablement 與 adoption 跟技術本身一樣重要。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.98,
        1.62,
        3.55,
        4.95,
        "一句話",
        [
            "我不是只會做系統，也能把系統講清楚、教清楚、交接清楚。這對 helixGPT 類產品其實很關鍵。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_gap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "9. 我要誠實講的 gap")
    add_text_panel(
        slide,
        0.75,
        1.62,
        3.9,
        4.95,
        "我不會硬講的事",
        [
            "我不會說我知道 helixGPT 的內部架構。",
            "我不會假裝自己已經用過所有企業 AI 平台工具。",
            "我不會把個人專案包裝成大公司上線案例。",
        ],
        fill=RED_SOFT,
    )
    add_text_panel(
        slide,
        4.9,
        1.62,
        3.8,
        4.95,
        "但我有的能力",
        [
            "AI 助手架構與 RAG 方法。",
            "資料工程、系統設計、維運與觀測。",
            "文件化、解釋能力、持續優化能力。",
        ],
        fill=GREEN_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.62,
        3.55,
        4.95,
        "我會怎麼補強",
        [
            "先快速理解產品定位與使用者工作流。",
            "用現有 agent / RAG / operations 經驗去對位內部系統。",
            "把自己已有能力接到團隊既有流程上，而不是從零開始摸索。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_306090(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "10. 30 / 60 / 90 天我會怎麼做")
    add_text_panel(
        slide,
        0.82,
        1.7,
        3.75,
        4.85,
        "前 30 天",
        [
            "理解產品定位、使用者工作流、主要知識來源與常見失敗案例。",
            "快速把自己現在的 RAG / agent 經驗對上團隊現況。",
            "先接住文件整理、問題分析與小型優化任務。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.82,
        1.7,
        3.75,
        4.85,
        "前 60 天",
        [
            "獨立處理一條 AI 助手或知識檢索流程的優化。",
            "補強 evaluation、維運與使用者導入材料。",
            "沉澱可複用的 troubleshooting pattern。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.82,
        1.7,
        3.7,
        4.85,
        "前 90 天",
        [
            "能獨立扛一個子模組：例如知識檢索、回答品質、維運監控或導入文件。",
            "開始對產品與導入流程提出具體優化建議。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_talk_track(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "11. 面試時建議講法")
    add_text_panel(
        slide,
        0.8,
        1.62,
        6.0,
        4.95,
        "90 秒版本",
        [
            "如果這個角色本質上是在做 helixGPT 類企業 AI 助手，我覺得我最有價值的地方，不是某個特定平台工具，而是我已經做過完整的 AI 助手閉環：資料整合、RAG、回答策略、維運、文件與持續優化。",
            "Job Radar 雖然是求職場景，但底層能力和企業知識助手很接近。這也是我認為自己能快速進入狀況的原因。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.05,
        1.62,
        5.45,
        4.95,
        "被追問時",
        [
            "問 agent：講 routing、retrieval、tool-like workflow。",
            "問導入：講需求、PoC、驗收、交付。",
            "問維運：講延遲、精準度、cache、monitoring。",
            "問你本人：講你怎麼定義問題、怎麼做驗證。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def slide_appendix(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Appendix. 建議一起準備的材料")
    add_bullet_list(
        slide,
        0.95,
        1.75,
        11.2,
        4.9,
        [
            "這份 helixGPT 類 AI agent 對位 deck。",
            "thesis-style interview deck 作為 deeper dive 備份。",
            "AI / LLM / RAG report 頁，面試官若偏技術可往這邊帶。",
            "完整結構圖頁與系統架構頁，當白板替代材料。",
        ],
        font_size=17,
    )
    add_footer(slide, "Job Radar · HelixGPT-style Agent Deck")


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_what_role_is(prs)
    slide_mapping(prs)
    slide_agent_arch(prs)
    slide_use_cases(prs)
    slide_delivery(prs)
    slide_ops(prs)
    slide_metrics(prs)
    slide_docs(prs)
    slide_gap(prs)
    slide_306090(prs)
    slide_talk_track(prs)
    slide_appendix(prs)
    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate Job Radar helixGPT-style agent interview deck.")
    parser.add_argument("--output", default=DEFAULT_OUTPUT, help="Output PPTX path")
    args = parser.parse_args()
    output = Path(args.output).expanduser()
    output.parent.mkdir(parents=True, exist_ok=True)
    build_presentation().save(output)
    print(output)


if __name__ == "__main__":
    main()
