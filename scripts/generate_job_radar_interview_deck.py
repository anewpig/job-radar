from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(247, 244, 238)
PANEL = RGBColor(255, 255, 255)
INK = RGBColor(26, 39, 52)
MUTED = RGBColor(92, 101, 117)
ACCENT = RGBColor(40, 88, 116)
ACCENT_SOFT = RGBColor(225, 235, 242)
GOLD = RGBColor(185, 145, 67)
GOLD_SOFT = RGBColor(243, 236, 220)
RED_SOFT = RGBColor(245, 228, 224)
GREEN_SOFT = RGBColor(226, 241, 234)
LINE = RGBColor(213, 219, 226)

FONT = "PingFang TC"
FONT_EN = "Aptos"


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(9.5), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = INK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.12), Inches(11.2), Inches(0.5))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT
        run2.font.size = Pt(11)
        run2.font.color.rgb = MUTED


def add_badge(slide, text: str, x: float, y: float, w: float = 1.35) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.33)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_SOFT
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_EN
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: list[str],
    fill: RGBColor = PANEL,
    title_color: RGBColor = INK,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.bold = True
    run.font.size = Pt(15)
    run.font.color.rgb = title_color

    for item in body:
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(11.5)
        run.font.color.rgb = INK


def add_bullet_list(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: list[str],
    font_size: float = 16,
    color: RGBColor = INK,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        p.bullet = True
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def add_kpi(slide, x: float, y: float, w: float, h: float, title: str, value: str, tint: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = tint
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT
    r1.font.size = Pt(11)
    r1.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = FONT_EN
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = INK


def add_flow_node(slide, x: float, y: float, w: float, h: float, title: str, subtitle: str, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(9)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = INK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = FONT
    r2.font.size = Pt(10)
    r2.font.color.rgb = MUTED


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


def add_footer(slide, text: str = "Job Radar Interview Deck") -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(7.03), Inches(12), Inches(0.22))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_EN
    run.font.size = Pt(8)
    run.font.color.rgb = MUTED


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.62), Inches(12), Inches(5.9)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PANEL
    band.line.color.rgb = LINE

    add_badge(slide, "INTERVIEW", 0.92, 0.9, 1.45)
    title_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.4), Inches(7.2), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Job Radar\n求職工作台面試簡報"
    r.font.name = FONT
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = INK

    sub = slide.shapes.add_textbox(Inches(0.98), Inches(3.0), Inches(6.2), Inches(1.35))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = (
        "用產品視角 + 系統設計視角，說明這個專案如何把「職缺爬蟲」變成"
        "「求職工作台」，並支撐自動刷新、履歷分析、AI 助理與通知流程。"
    )
    r2.font.name = FONT
    r2.font.size = Pt(15)
    r2.font.color.rgb = MUTED

    add_kpi(slide, 7.85, 1.45, 2.0, 1.25, "核心主題", "Product", GOLD_SOFT)
    add_kpi(slide, 9.95, 1.45, 2.0, 1.25, "核心主題", "System", ACCENT_SOFT)
    add_kpi(slide, 7.85, 2.9, 2.0, 1.25, "核心主題", "AI", GREEN_SOFT)
    add_kpi(slide, 9.95, 2.9, 2.0, 1.25, "核心主題", "Data", RED_SOFT)

    add_text_panel(
        slide,
        7.85,
        4.45,
        4.1,
        1.55,
        "這份 deck 適合怎麼講",
        [
            "前 2 分鐘講產品問題與定位",
            "中間 6 到 8 分鐘講架構、資料流、AI 與 DB",
            "最後 2 分鐘講 trade-off、擴充與你接下來會怎麼做",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide)


def slide_product_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "1. 這個產品在解什麼問題", "先講使用者痛點，再講你為什麼把它做成產品，而不是單純爬蟲。")
    add_text_panel(
        slide,
        0.75,
        1.55,
        3.85,
        4.8,
        "使用者痛點",
        [
            "職缺分散在 104 / 1111 / Cake / LinkedIn，求職者需要重複查找。",
            "每次回來都要重新搜尋、重新判斷哪些值得投。",
            "履歷與市場需求之間的落差很難量化，容易只靠感覺投遞。",
            "通知、追蹤、收藏、投遞看板若分散在多處，求職流程很難形成閉環。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.83,
        1.55,
        3.85,
        4.8,
        "產品定義",
        [
            "不是一支 script，而是求職工作台。",
            "核心價值鏈是：搜尋 → 快照 → 分析 → 決策 → 追蹤 → 通知 → 再回來。",
            "把一次性的查詢工具，變成可持續回訪的個人化系統。",
            "AI 不當 source of truth，而是做分析、摘要、問答與建議。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.92,
        1.55,
        3.65,
        4.8,
        "面試時要強調",
        [
            "我先定義產品邊界，再決定工程切法。",
            "這套系統的重點不是抓多少站，而是如何把資料轉成持續可用的決策介面。",
            "從 day-1 就有同步 / 非同步、快取 / 歷史、產品資料 / 執行期資料的分層觀念。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide)


def slide_jtbd(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "2. 產品定位、JTBD 與回訪循環")
    add_flow_node(slide, 0.8, 2.05, 2.05, 1.0, "使用者輸入條件", "角色 / 地區 / 關鍵字 / 頻率", PANEL)
    add_flow_node(slide, 3.1, 2.05, 2.05, 1.0, "系統抓取快照", "normalize + dedupe + enrich", ACCENT_SOFT)
    add_flow_node(slide, 5.4, 2.05, 2.05, 1.0, "分析與排序", "相關度 / 技能 / 任務 insights", GREEN_SOFT)
    add_flow_node(slide, 7.7, 2.05, 2.05, 1.0, "做投遞決策", "收藏 / 看板 / 履歷匹配", GOLD_SOFT)
    add_flow_node(slide, 10.0, 2.05, 2.05, 1.0, "回訪與通知", "saved search + refresh + alerts", RED_SOFT)
    for start in [2.85, 5.15, 7.45, 9.75]:
        add_arrow(slide, start, 2.55, start + 0.25, 2.55)

    add_text_panel(
        slide,
        0.85,
        3.55,
        6.0,
        2.35,
        "主要 JTBD",
        [
            "當我正在找工作時，我想把分散的市場資訊收斂成固定的工作台，這樣我可以持續追蹤機會，不會漏掉重要職缺。",
            "當我不確定自己是否符合某些職缺時，我希望系統能把履歷與市場需求對照，幫我看到缺口與優先補強方向。",
            "當我回訪系統時，我不想從零開始，而是看到最新變化、通知與該處理的下一步。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.05,
        3.55,
        5.4,
        2.35,
        "為什麼這樣設計",
        [
            "這個定位會自然推導出 saved search、snapshot、history、notifications、tracking center。",
            "如果只把它當爬蟲，系統就會停在一次性抓資料；如果把它當工作台，就必須處理狀態、回訪、持續更新與長期資料。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide)


def slide_topology(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "3. 整體系統拓樸")
    layers = [
        ("使用者互動層", "Streamlit UI / pages / navigation / launcher", ACCENT_SOFT, 0.9),
        ("UI 協調層", "bootstrap / session / router / search setup / crawl runtime", PANEL, 3.25),
        ("核心服務層", "crawl application service / pipeline / analysis / notifications", GOLD_SOFT, 5.6),
        ("執行與持久化", "scheduler / worker / SQLite DBs / snapshots / history", GREEN_SOFT, 7.95),
        ("外部系統", "104 / 1111 / Cake / LinkedIn / LINE / Email / LLM", RED_SOFT, 10.3),
    ]
    for title, desc, fill, x in layers:
        add_flow_node(slide, x, 2.0, 2.0, 2.45, title, desc, fill)
    for start in [2.95, 5.3, 7.65, 10.0]:
        add_arrow(slide, start, 3.2, start + 0.22, 3.2)

    add_text_panel(
        slide,
        0.9,
        5.05,
        12.0,
        1.55,
        "我會怎麼向面試官解釋",
        [
            "這是一個模組化單體，不是微服務。原因是功能面很多，但團隊規模與 stage 還沒到要付出微服務協調成本。",
            "我把邏輯拆成明確的分層與多進程：UI 負責互動；scheduler 只決定該不該做；worker 真正執行抓取與同步；資料則依產品主資料、執行期資料、歷史 archive 分庫。",
        ],
        fill=PANEL,
    )
    add_footer(slide)


def slide_frontend(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "4. 前端架構：不是只有畫面，而是工作台 orchestration")
    add_text_panel(
        slide,
        0.75,
        1.6,
        3.0,
        4.95,
        "App Shell",
        [
            "bootstrap 先建立 service wiring 與 user context。",
            "session state 管 active user、main tab、saved search、crawl pending 狀態。",
            "router / navigation 決定不同頁面的 surface 與切換。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        3.95,
        1.6,
        3.0,
        4.95,
        "搜尋工作台",
        [
            "search setup 管多組搜尋條件、頻率、標籤與輸入列。",
            "crawl runtime 負責啟動查詢、顯示狀態、輪詢 background job、收斂結果。",
            "overview / market / tracking / board 組成使用者的決策與回訪介面。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        7.15,
        1.6,
        2.75,
        4.95,
        "AI 與輔助入口",
        [
            "resume page 做履歷匹配。",
            "assistant page 做較完整的 AI 助理。",
            "floating launcher 提供輕量問答、說明與通知入口。",
        ],
        fill=GOLD_SOFT,
    )
    add_text_panel(
        slide,
        10.1,
        1.6,
        2.45,
        4.95,
        "這樣做的優點",
        [
            "功能很多，但頁面責任相對清楚。",
            "session state 讓 Streamlit 也能支撐多步驟互動。",
            "透過後續重構，把大檔拆成 orchestration / sections / actions，比較能維持可讀性。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide)


def slide_backend(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "5. 後端執行架構：同步入口 + 非同步處理")
    add_flow_node(slide, 0.9, 2.2, 2.0, 0.95, "App", "接住使用者請求，做最小同步工作", PANEL)
    add_flow_node(slide, 3.3, 2.2, 2.1, 0.95, "Crawl App Service", "cache / schedule / submit / sync", ACCENT_SOFT)
    add_flow_node(slide, 5.8, 2.2, 1.95, 0.95, "Runtime Queue", "job / lease / signal / finalize", GOLD_SOFT)
    add_flow_node(slide, 8.1, 2.2, 1.85, 0.95, "Worker", "fetch / analyze / notify", GREEN_SOFT)
    add_flow_node(slide, 10.3, 2.2, 1.95, 0.95, "Scheduler", "refresh due saved searches", RED_SOFT)
    add_arrow(slide, 2.95, 2.67, 3.18, 2.67)
    add_arrow(slide, 5.48, 2.67, 5.62, 2.67)
    add_arrow(slide, 7.84, 2.67, 8.02, 2.67)
    add_arrow(slide, 10.05, 2.67, 9.97, 2.67)

    add_text_panel(
        slide,
        0.95,
        3.75,
        12.0,
        2.25,
        "為什麼這樣切",
        [
            "使用者第一次查詢時，系統可以同步命中快取，也可以把真正耗時的抓取交給 background worker，避免 UI 卡死。",
            "scheduler 只負責找出 due saved searches，不直接抓資料；worker 只負責執行，這樣觀察性與維運成本比較低。",
            "queue 與 runtime signal 獨立出來後，前端可以用輪詢看進度，而不是把整段長任務塞在 Streamlit request 裡。",
        ],
        fill=PANEL,
    )
    add_footer(slide)


def slide_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "6. 抓取管線：Connector → Normalize → Analyze → Snapshot")
    stages = [
        ("Connector", "各站搜尋 / detail 抽取", PANEL),
        ("Normalize", "欄位映射、時間與薪資正規化", ACCENT_SOFT),
        ("Dedupe", "避免重複職缺污染快照", GOLD_SOFT),
        ("Analyze", "角色匹配、技能、任務 insights", GREEN_SOFT),
        ("Persist", "market snapshot / runtime / history", RED_SOFT),
    ]
    x = 0.8
    for title, subtitle, fill in stages:
        add_flow_node(slide, x, 2.0, 2.25, 1.15, title, subtitle, fill)
        x += 2.45
    for start in [3.02, 5.47, 7.92, 10.37]:
        add_arrow(slide, start, 2.57, start + 0.18, 2.57)

    add_text_panel(
        slide,
        0.85,
        3.75,
        5.9,
        2.4,
        "優點",
        [
            "來源差異被收斂在 connector 與 normalize 層。",
            "後面的分析與產品頁面可以面向統一資料模型，不必處理各站特例。",
            "快照與歷史都能吃同一批結構化結果。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.0,
        3.75,
        5.55,
        2.4,
        "缺點與替代方案",
        [
            "現在還是模組化單體，connector、analysis、persist 在同一個部署邊界內。",
            "若來源量暴增，可以把抓取層進一步拆成獨立 ingestion service，或改成事件流式架構。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide)


def slide_ai(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "7. AI 架構：不是把模型塞進產品，而是放在正確的位置")
    add_text_panel(
        slide,
        0.75,
        1.6,
        3.85,
        4.85,
        "履歷分析",
        [
            "先做文字抽取與 rule-based profile 整理。",
            "再用 LLM 補強職能理解、技能聚合與摘要。",
            "失敗時可 fallback，不讓產品功能整段失效。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.8,
        1.6,
        3.85,
        4.85,
        "履歷匹配與建議",
        [
            "先用規則與結構化欄位做全量篩選。",
            "再對候選職缺做更重的語意比對與打分。",
            "避免把所有職缺都丟給模型，控制成本與延遲。",
        ],
        fill=GREEN_SOFT,
    )
    add_text_panel(
        slide,
        8.85,
        1.6,
        3.7,
        4.85,
        "AI 助理 / RAG",
        [
            "以當前 snapshot 為中心做 grounded answer。",
            "不是直接吃整個歷史資料庫，避免答案飄掉。",
            "同時保留 monitoring、blocked state 與 fallback response。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide)


def slide_db(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "8. DB 與持久化：多庫不是複雜，而是邊界管理")
    add_text_panel(
        slide,
        0.75,
        1.65,
        2.95,
        4.8,
        "product_state.sqlite3",
        [
            "使用者偏好、saved search、通知設定、投遞看板等產品主資料。",
            "這層最接近使用者心智模型，也最該穩定。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        3.95,
        1.65,
        2.95,
        4.8,
        "query_runtime.sqlite3",
        [
            "查詢 queue、job lease、runtime signal、snapshot runtime 狀態。",
            "這些是執行期資料，不應污染產品主資料。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        7.15,
        1.65,
        2.75,
        4.8,
        "user_submissions.sqlite3",
        [
            "履歷文本、表單提交與較敏感的使用者上傳資料。",
            "分開存放有利於資料治理與後續加密策略。",
        ],
        fill=GOLD_SOFT,
    )
    add_text_panel(
        slide,
        10.15,
        1.65,
        2.35,
        4.8,
        "market_history.sqlite3",
        [
            "歷史市場 archive。",
            "未來量大時可以先搬這層到分析型儲存，而不用整站一起搬。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide)


def slide_data_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "9. 四條最重要的資料流")
    flows = [
        ("手動查詢", "使用者輸入條件 → start crawl → worker 跑 pipeline → snapshot 回 UI", 0.9, PANEL),
        ("自動刷新", "saved search 到期 → scheduler enqueue → worker finalize → 通知 / tracking 更新", 3.65, ACCENT_SOFT),
        ("履歷分析", "上傳履歷 → extract text → AI / rule 分析 → profile / match result 落地", 6.4, GREEN_SOFT),
        ("AI 問答", "輸入問題 → 檢索 snapshot context → RAG answer / fallback → 顯示結果", 9.15, GOLD_SOFT),
    ]
    for title, subtitle, x, fill in flows:
        add_flow_node(slide, x, 2.0, 2.35, 1.25, title, subtitle, fill)
    add_text_panel(
        slide,
        0.95,
        3.85,
        11.7,
        2.1,
        "我在面試會怎麼講這一頁",
        [
            "我通常不會逐頁背模組名稱，而是抓四條資料流講。因為資料流最能讓面試官判斷你是不是真的理解系統。",
            "這四條流分別代表：產品主流程、背景刷新、AI 加值、以及互動式問答。只要這四條你講得清楚，整體設計就站得住。",
        ],
        fill=PANEL,
    )
    add_footer(slide)


def slide_tradeoffs(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "10. 為什麼現在這樣做：優點、缺點、替代方案")
    add_text_panel(
        slide,
        0.75,
        1.65,
        3.9,
        4.75,
        "目前方案",
        [
            "模組化單體 + Streamlit UI",
            "多 SQLite 分庫",
            "scheduler / worker 多進程",
            "snapshot-centered AI",
            "connector + pipeline 統一資料模型",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.9,
        1.65,
        3.55,
        4.75,
        "優點",
        [
            "開發快、部署簡單、邏輯集中，適合快速驗證產品。",
            "能先做正確的資料與狀態邊界，而不是一開始就被分散式複雜度壓垮。",
            "對 side project 或小團隊來說，性價比高。",
        ],
        fill=GREEN_SOFT,
    )
    add_text_panel(
        slide,
        8.7,
        1.65,
        3.85,
        4.75,
        "缺點 / 替代",
        [
            "高併發與多人協作擴張時，單體會逐步逼近上限。",
            "替代方案可以是前後端分離、Postgres + queue、獨立 ingestion service、事件流架構。",
            "但這些要在產品已被驗證後再做，才有投資報酬率。",
        ],
        fill=RED_SOFT,
    )
    add_footer(slide)


def slide_scaling(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "11. 如果要擴充，我會怎麼演進")
    add_flow_node(slide, 0.95, 2.15, 2.6, 1.05, "Stage 1：現在", "模組化單體 + worker + SQLite", PANEL)
    add_flow_node(slide, 3.85, 2.15, 2.6, 1.05, "Stage 2：產品成長", "前後端分離 / Postgres / Redis queue", ACCENT_SOFT)
    add_flow_node(slide, 6.75, 2.15, 2.6, 1.05, "Stage 3：資料規模化", "獨立 ingestion / warehouse / analytics", GOLD_SOFT)
    add_flow_node(slide, 9.65, 2.15, 2.6, 1.05, "Stage 4：企業化", "多租戶 / 權限 / audit /治理", GREEN_SOFT)
    add_arrow(slide, 3.58, 2.67, 3.75, 2.67)
    add_arrow(slide, 6.48, 2.67, 6.65, 2.67)
    add_arrow(slide, 9.38, 2.67, 9.55, 2.67)

    add_text_panel(
        slide,
        1.0,
        3.8,
        11.5,
        2.15,
        "擴充原則",
        [
            "先搬最痛的層，不是一次全搬。通常第一個搬的是 queue 與主資料庫，第二個搬的是 ingestion，最後才是把單體拆成多服務。",
            "AI 也一樣：先補 cache、cost control、eval，再談更大模型或多 agent。否則成本跟複雜度會先失控。",
            "資料層要維持 ownership：產品主資料、執行期資料、歷史 archive、敏感提交，這四層不要重新混在一起。",
        ],
        fill=PANEL,
    )
    add_footer(slide)


def slide_interview_talk(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "12. 面試時我會怎麼講")
    add_text_panel(
        slide,
        0.8,
        1.65,
        6.1,
        4.85,
        "3 分鐘版本",
        [
            "這個專案表面上是職缺聚合，實際上我把它做成求職工作台。",
            "我先定義使用者回訪循環，所以系統自然長出 saved search、snapshot、tracking、notifications 與 AI 分析。",
            "在架構上我採模組化單體，但把 UI、scheduler、worker、runtime queue 與多個 SQLite 邊界切清楚，讓它能從 side project 演進到產品。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.15,
        1.65,
        5.4,
        4.85,
        "面試官最可能追問",
        [
            "為什麼不用微服務？",
            "為什麼不用單一 DB？",
            "AI 放在哪些位置，怎麼避免亂答？",
            "如果來源、使用者、通知量暴增，你先擴哪裡？",
            "你怎麼知道這是一個產品，而不是一堆功能？",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar Interview Deck · 建議配合你自己的口頭版本做第二次微調")


def slide_appendix(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Appendix. 你可以順手附給面試官的文件")
    add_bullet_list(
        slide,
        0.95,
        1.7,
        11.4,
        4.8,
        [
            "總架構手冊：job-radar-system-architecture.html",
            "完整結構圖：job-radar-complete-structure-diagram.html",
            "AI 詳細流程：job-radar-ai-detailed-flow.html",
            "DB Schema 手冊：job-radar-db-schema-handbook.html",
            "產品與系統 Playbook：job-radar-product-system-playbook.html",
            "學習地圖：job-radar-learning-curriculum.html",
            "如果要更像面試附件，可以把這份 PPT 搭配 1 頁 executive summary 一起給。",
        ],
        font_size=17,
    )
    add_footer(slide)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_product_problem(prs)
    slide_jtbd(prs)
    slide_topology(prs)
    slide_frontend(prs)
    slide_backend(prs)
    slide_pipeline(prs)
    slide_ai(prs)
    slide_db(prs)
    slide_data_flow(prs)
    slide_tradeoffs(prs)
    slide_scaling(prs)
    slide_interview_talk(prs)
    slide_appendix(prs)
    return prs


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate the Job Radar interview PowerPoint deck.")
    parser.add_argument(
        "--output",
        default="/Users/zhuangcaizhen/Desktop/專案/UI練習/job-radar-interview-deck.pptx",
        help="Output .pptx path",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(output))
    print(output)


if __name__ == "__main__":
    main()
