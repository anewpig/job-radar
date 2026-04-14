from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation

from generate_job_radar_interview_deck import (
    ACCENT_SOFT,
    BG,
    GOLD_SOFT,
    GREEN_SOFT,
    PANEL,
    RED_SOFT,
    SLIDE_H,
    SLIDE_W,
    add_arrow,
    add_badge,
    add_bullet_list,
    add_footer,
    add_flow_node,
    add_kpi,
    add_text_panel,
    add_title,
    set_background,
)


ROLE_TITLE = "GCP No-Code / Low-Code AI 導入顧問向"
DEFAULT_OUTPUT = "/Users/zhuangcaizhen/Desktop/專案/UI練習/job-radar-gcp-ai-role-deck.pptx"


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    band = slide.shapes.add_shape(1, 0.7 * 914400, 0.62 * 914400, 12 * 914400, 5.95 * 914400)
    band.fill.solid()
    band.fill.fore_color.rgb = PANEL
    band.line.color.rgb = BG

    add_badge(slide, "ROLE FIT", 0.95, 0.92, 1.4)
    add_title(
        slide,
        "Job Radar 專案\n對應 GCP AI 導入職缺的面試簡報",
        "把現有專案改寫成職缺對位版，重點講 AI 導入、RAG / Agent 能力、平台維運與客戶導入思維。",
    )
    add_kpi(slide, 8.0, 1.65, 1.95, 1.22, "關鍵能力", "AI", GREEN_SOFT)
    add_kpi(slide, 10.15, 1.65, 1.95, 1.22, "關鍵能力", "RAG", ACCENT_SOFT)
    add_kpi(slide, 8.0, 3.0, 1.95, 1.22, "關鍵能力", "Agent", GOLD_SOFT)
    add_kpi(slide, 10.15, 3.0, 1.95, 1.22, "關鍵能力", "Ops", RED_SOFT)
    add_text_panel(
        slide,
        0.95,
        3.05,
        6.15,
        2.15,
        "這份 deck 的目標",
        [
            "不是單純介紹專案，而是向面試官證明：我已具備把 AI 功能做成產品、把資料流程做成可運營系統、把技術內容講成導入方案的能力。",
            "同時誠實區分：哪些是我已做過的，哪些是可以直接遷移到 GCP stack 的，哪些需要上線前再補實戰驗證。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_job_understanding(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "1. 我怎麼理解這個職缺", "這份工作不是純模型研究，而是把 AI 工具導入客戶工作流程。")
    add_text_panel(
        slide,
        0.75,
        1.6,
        3.95,
        4.95,
        "職位核心",
        [
            "以 GCP No-Code / Low-Code AI 服務為核心，快速把 AI 解法導入客戶場景。",
            "重點不是從零訓練模型，而是把現成 AI 能力組成可落地的流程、助手、分析工具與導入方案。",
            "角色同時涵蓋 solution design、部署、客戶需求分析、維運與技術溝通。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.95,
        1.6,
        3.8,
        4.95,
        "JD 關鍵能力",
        [
            "GCP AI 產品熟悉度：Vertex AI、Agent Builder、Gemini、Workspace AI。",
            "生成式 AI 應用：LLM、RAG、Agent。",
            "平台操作：gcloud、Cloud Console、Linux。",
            "客戶導入：需求分析、技術支援、溝通、文件化。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.6,
        3.55,
        4.95,
        "面試時真正要證明",
        [
            "我不是只會講 AI 名詞，而是能把需求拆成導入方案。",
            "我可以把產品與系統做成可維運的形態。",
            "我能向非技術利害關係人解釋價值，也能向技術方解釋架構與風險。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_project_reframe(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "2. 我會怎麼把 Job Radar 專案重新包裝成這個職缺想看的案例")
    add_flow_node(slide, 0.9, 2.0, 2.2, 1.0, "問題定義", "求職市場資訊分散、分析成本高", PANEL)
    add_flow_node(slide, 3.45, 2.0, 2.2, 1.0, "AI 導入設計", "搜尋 / 分析 / 問答 / 推薦", ACCENT_SOFT)
    add_flow_node(slide, 6.0, 2.0, 2.2, 1.0, "系統落地", "crawler / snapshot / RAG / notification", GREEN_SOFT)
    add_flow_node(slide, 8.55, 2.0, 2.2, 1.0, "使用者工作流", "回訪、決策、履歷補強", GOLD_SOFT)
    add_flow_node(slide, 11.1, 2.0, 1.4, 1.0, "營運", "refresh / monitor", RED_SOFT)
    for start in [3.15, 5.7, 8.25, 10.8]:
        add_arrow(slide, start, 2.5, start + 0.18, 2.5)
    add_text_panel(
        slide,
        0.95,
        3.6,
        11.5,
        2.35,
        "包裝重點",
        [
            "我不會把它講成『我做了一個職缺爬蟲』，而是講成『我把分散資訊場景變成一個 AI-enhanced 工作台，從資料收集、分析、問答到通知形成完整閉環』。",
            "這正好對應 AI 導入角色的核心能力：理解場景、設計流程、選技術、做導入、持續維運。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_fit_matrix(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "3. JD 對位矩陣：我已證明的能力 vs 可轉移能力")
    add_text_panel(
        slide,
        0.72,
        1.6,
        3.05,
        4.95,
        "已證明",
        [
            "生成式 AI 應用：LLM、RAG、AI 助理、履歷分析。",
            "多來源資料導入與正規化。",
            "把 AI 功能做成可使用、可維運的產品流程。",
            "技術文件、架構頁、流程圖、簡報化表達。",
        ],
        fill=GREEN_SOFT,
    )
    add_text_panel(
        slide,
        4.0,
        1.6,
        4.1,
        4.95,
        "可直接遷移到 JD",
        [
            "把 snapshot-centered RAG 遷移到 Vertex AI Search / Agent Builder 思維。",
            "把目前 AI orchestration 映射成 Gemini + Workspace + API integration 流程。",
            "把目前 scheduler / worker / monitor 概念對應到雲端導入後的運營與維運責任。",
            "把產品需求拆解、風險分析、文件化，轉成客戶導入流程。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.35,
        1.6,
        4.2,
        4.95,
        "需要誠實補強",
        [
            "目前專案不是直接跑在 GCP Vertex / Agent Builder 上。",
            "gcloud / Cloud Console 的實作深度，若要面對 production 客戶環境，還要再補更完整案例。",
            "這不是弱點，重點是我知道哪些能力已具備、哪些要在前 30 到 60 天補齊。",
        ],
        fill=RED_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_ai_mapping(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "4. 生成式 AI 能力：這個專案和 JD 最強的交集")
    add_text_panel(
        slide,
        0.75,
        1.58,
        3.9,
        5.0,
        "LLM / RAG / Agent",
        [
            "專案內已有 AI 助理、履歷分析、問答與任務/技能萃取流程。",
            "RAG 不是空談，而是建立在 snapshot 與檢索 context 上。",
            "這讓我在面試時可以直接講 retrieval、grounding、fallback、monitoring，而不是只背工具名。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.9,
        1.58,
        3.75,
        5.0,
        "對應 GCP 產品語言",
        [
            "RAG / grounded answer 可映射到 Vertex AI Search / Agent Builder / Gemini workflow。",
            "對話式 AI 能力可對應 Conversation / Agent Builder 的設計思路。",
            "生成式摘要、建議、分析可對應 Gemini for Google Cloud 與 Workspace AI 協作場景。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.9,
        1.58,
        3.6,
        5.0,
        "面試時建議講法",
        [
            "我有做過把資料流、檢索、回答策略包成實際產品功能的經驗。",
            "如果改用 GCP stack，我會把既有流程拆成：資料入口、檢索層、Gemini 推理層、agent orchestration、觀測層。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_gcp_mapping(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "5. 如果把這個專案搬到 GCP，我會怎麼映射")
    add_flow_node(slide, 0.85, 2.0, 2.2, 1.0, "資料入口", "crawler / API / 表單 / 上傳", PANEL)
    add_flow_node(slide, 3.35, 2.0, 2.2, 1.0, "資料與快照", "Cloud Storage / DB / index", ACCENT_SOFT)
    add_flow_node(slide, 5.85, 2.0, 2.2, 1.0, "AI 能力層", "Gemini / Vertex AI / Agent Builder", GREEN_SOFT)
    add_flow_node(slide, 8.35, 2.0, 2.2, 1.0, "工作流層", "客服 / 助手 / 分析 / 協作", GOLD_SOFT)
    add_flow_node(slide, 10.85, 2.0, 1.65, 1.0, "維運層", "logging / monitor / IAM", RED_SOFT)
    for start in [3.05, 5.55, 8.05, 10.55]:
        add_arrow(slide, start, 2.5, start + 0.18, 2.5)
    add_text_panel(
        slide,
        0.95,
        3.7,
        11.45,
        2.3,
        "重點不是逐字對產品名，而是逐層映射",
        [
            "我會跟面試官說：我已經有把 AI 功能做成產品與資料流程的經驗，若進公司後要改用 GCP 服務，我會先保留邏輯分層，再把每一層映射到合適的 managed service。",
            "這代表我理解的是系統問題，而不是只會使用某一個 UI 工具。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_delivery(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "6. 客戶導入與需求分析：我會怎麼把技術做成導入專案")
    add_flow_node(slide, 1.0, 2.1, 2.2, 0.95, "需求訪談", "場景、資料、流程、限制", PANEL)
    add_flow_node(slide, 3.6, 2.1, 2.2, 0.95, "PoC 定義", "成功指標、資料邊界、風險", ACCENT_SOFT)
    add_flow_node(slide, 6.2, 2.1, 2.2, 0.95, "導入實作", "AI 流程、介面、權限、維運", GREEN_SOFT)
    add_flow_node(slide, 8.8, 2.1, 2.2, 0.95, "驗收與文件", "操作說明、問題案例、交接", GOLD_SOFT)
    add_arrow(slide, 3.3, 2.55, 3.48, 2.55)
    add_arrow(slide, 5.9, 2.55, 6.08, 2.55)
    add_arrow(slide, 8.5, 2.55, 8.68, 2.55)
    add_text_panel(
        slide,
        1.0,
        3.75,
        11.2,
        2.25,
        "為什麼我可以講這一塊",
        [
            "因為這個專案本身就包含產品定義、資料流、AI 功能、維運機制與文件化輸出。我不是只做 demo，而是把整個流程從需求、實作到運營走過一次。",
            "這跟 AI 導入角色很接近：你要把技術能力轉成客戶能接受、能操作、能長期使用的方案。",
        ],
        fill=PANEL,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_ops_case(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "7. 維運與問題排除：我會拿什麼案例來說")
    add_text_panel(
        slide,
        0.75,
        1.6,
        3.85,
        5.05,
        "案例一：長任務與 UI 解耦",
        [
            "問題：抓取、分析、通知若都塞在同一個互動請求，UI 會卡住也難觀測。",
            "做法：切成 app / scheduler / worker / runtime queue。",
            "結果：同步入口更輕、背景處理更穩、前端能顯示進度。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.85,
        1.6,
        3.85,
        5.05,
        "案例二：資料邊界與可維運性",
        [
            "問題：產品主資料、執行期資料、歷史資料、敏感提交若混在一起，後續很難治理。",
            "做法：拆成多個 SQLite，讓 ownership 清楚。",
            "結果：備份、搬移、治理、擴充路線更清楚。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.6,
        3.55,
        5.05,
        "案例三：AI 可靠性",
        [
            "問題：AI 功能容易因上下文不足或模型回應不穩而失效。",
            "做法：snapshot-centered context、monitoring、fallback-first。",
            "結果：把 AI 從 demo 變成可交付功能。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_gap_plan(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "8. 差距我會怎麼說：不硬拗，但有補強方案")
    add_text_panel(
        slide,
        0.75,
        1.65,
        3.9,
        4.9,
        "目前不是的事",
        [
            "這個專案目前不是直接用 Vertex AI / Agent Builder 上線。",
            "不是大型企業客戶的正式 GCP 導入專案。",
            "不是純 no-code 操作型案例，而是系統與產品設計更強。",
        ],
        fill=RED_SOFT,
    )
    add_text_panel(
        slide,
        4.9,
        1.65,
        3.8,
        4.9,
        "但我已經具備",
        [
            "AI 流程設計能力",
            "RAG / Agent / 資料流與產品整合能力",
            "文件化、架構說明、簡報與溝通能力",
            "把功能做成可維運系統的能力",
        ],
        fill=GREEN_SOFT,
    )
    add_text_panel(
        slide,
        8.95,
        1.65,
        3.55,
        4.9,
        "補強策略",
        [
            "前 30 天快速把 GCP 對應服務實作成 demo。",
            "補強 gcloud / IAM / 監控 / deployment runbook。",
            "把現有專案的一條 AI flow 映射成標準 GCP 導入範例。",
        ],
        fill=GOLD_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_306090(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "9. 如果我錄取，30 / 60 / 90 天我會怎麼做")
    add_text_panel(
        slide,
        0.82,
        1.7,
        3.75,
        4.85,
        "前 30 天",
        [
            "熟悉公司現有 GCP AI 產品與典型客戶場景。",
            "把現有 Job Radar 專案中的一條 AI flow 重新用 GCP stack 畫成 mapping。",
            "補齊 gcloud / IAM / deployment / monitor 操作細節。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        4.82,
        1.7,
        3.75,
        4.85,
        "前 60 天",
        [
            "跟著真實導入案走需求澄清、PoC 與驗收。",
            "整理可複用的導入模板：需求盤點、風險清單、技術文件、FAQ。",
            "把生成式 AI 能力與 Workspace / Agent Builder 場景串起來。",
        ],
        fill=ACCENT_SOFT,
    )
    add_text_panel(
        slide,
        8.82,
        1.7,
        3.7,
        4.85,
        "前 90 天",
        [
            "能獨立帶中小型導入案或獨立負責特定模組。",
            "把一套從需求分析到上線後維運的流程標準化。",
            "開始沉澱案例、最佳實務與可複用素材。",
        ],
        fill=GREEN_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_talk_track(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "10. 面試時建議講法")
    add_text_panel(
        slide,
        0.8,
        1.62,
        6.0,
        4.95,
        "90 秒版本",
        [
            "我現在這個專案雖然表面上是職缺平台，但本質上是把資料收集、AI 分析、問答、通知與使用者工作流做成可維運產品。",
            "這跟你們要找的 AI 導入角色很接近，因為核心能力都是：理解場景、選技術、把 AI 放進流程、做部署與長期運營。",
            "我已經做過 LLM、RAG、Agent-like workflow、monitoring 與資料流設計；如果進公司，我會把這些能力映射到 GCP stack，快速補上平台細節與導入實戰。",
        ],
        fill=PANEL,
    )
    add_text_panel(
        slide,
        7.05,
        1.62,
        5.45,
        4.95,
        "面試官追問時可延伸",
        [
            "請我畫資料流：我可以畫 app / worker / AI / DB 邊界。",
            "問我 AI 怎麼避免亂答：我講 grounded context、fallback、monitoring。",
            "問我怎麼導入客戶：我講需求分析、PoC、驗收、文件與維運。",
            "問我 GCP 經驗：我誠實講現況，再講映射策略與補強計畫。",
        ],
        fill=ACCENT_SOFT,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def slide_appendix(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Appendix. 我建議你面試時順手準備的補充材料")
    add_bullet_list(
        slide,
        0.95,
        1.75,
        11.2,
        4.9,
        [
            "1 頁專案摘要：一句產品定位 + 一張系統圖 + 三個關鍵能力。",
            "這份 GCP AI 角色對位版 deck。",
            "原本的通用版系統 deck，作為 deeper dive 備份。",
            "架構 HTML：總架構、AI 詳細流程、DB schema、完整結構圖。",
            "若可以，再補 1 個你會如何把這個專案遷到 GCP 的白板版流程。",
        ],
        font_size=17,
    )
    add_footer(slide, "Job Radar · GCP AI Role Deck")


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_job_understanding(prs)
    slide_project_reframe(prs)
    slide_fit_matrix(prs)
    slide_ai_mapping(prs)
    slide_gcp_mapping(prs)
    slide_delivery(prs)
    slide_ops_case(prs)
    slide_gap_plan(prs)
    slide_306090(prs)
    slide_talk_track(prs)
    slide_appendix(prs)
    return prs


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate GCP AI role interview deck for Job Radar.")
    parser.add_argument("--output", default=DEFAULT_OUTPUT, help="Output .pptx path")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
